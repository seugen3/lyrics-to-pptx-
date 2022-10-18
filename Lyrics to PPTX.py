import os
import time
from turtle import width
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import ColorFormat, RGBColor
from selenium import webdriver
from selenium.webdriver import ActionChains
from selenium.webdriver.common.by import By
from selenium.common.exceptions import TimeoutException
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
import asyncio

# Get all the pats
driverPath = input("Please input your webdriver path \n")
urlPath  =input("Please input here the URL \n")
titlePath = input("Insert the title xPath \n")
lyricsPath = input("Insert the xPath of lyrics \n")
folderPath = input("Enter the folder path for the saved file \n")

# Get Chrome driver
driver = webdriver.Chrome(driverPath)

# Get the URL
url = urlPath

# Acces the URL
driver.get(url)
    
# Find the title
title = driver.find_element('xpath', titlePath).text
print("Found the title!")

#Find the lyrics
lyrics = driver.find_element('xpath', lyricsPath).text
print("Found the lyrics!")

title_of_file = str(title) + ".pptx"
savedPath = str(folderPath) + str(title) + ".pptx"


# Develop the pptx File
prs = Presentation()
# Type of slide. You can change the "6" to see other types of slides
blank_slide_layout = prs.slide_layouts[6]


# Find how many strophes there are
n_Strophes = lyrics.split('\n\n')
nrStrophes = len(n_Strophes)

# Set the number of slides to -1. 
nrSlide = -1
print("I have the lyrics ready to be pasted!")

# This function crrate sets the number of slides.
def path_number():
    global nrSlide
    nrSlide += 1

# This functions creats the slides
def slide_chain():
    slide = prs.slides.add_slide(blank_slide_layout)
    # Set the backround as black
    background = slide.background
    fill = background.fill
    fill.solid() #Sets the backround as a solid color
    fill.background_color= (0, 0, 0) #This is the RGB color. You can try new colors
    # Set the span of text box. This setting makes the text box fil the entire slide
    width = Inches(16)
    left = Inches(0)
    height = Inches(9)
    top = Inches(0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    global tf
    tf = txBox.text_frame
    # Slide's shape. I preffer it 16:9
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

# Writting of lyrics on slides
def paragraph_chain():
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER # Allig the lyrics to center
    p.text = n_Strophes[nrSlide] # Each strophes to each slide
    p.font.name = 'Bahnschrift' # Set the fot. I like Bahnschrift
    p.font.size = Pt(55) # Font size
    p.font.bold = True # Bold or not
    p.font.color.rgb = RGBColor(255,255,255) # Font color
    print("I wrote the ", nrSlide+1, "verse")


# Main
while nrSlide < nrStrophes-1:
    nrSlide += 1
    if nrSlide < 1:
        print("I am writting the PowerPoint Presentation!")
    slide_chain()
    paragraph_chain()
    prs.save(savedPath)
    if nrSlide == nrStrophes-1:
        print("I saved the presentation! It's name is: ", title_of_file, "\n\n\n")